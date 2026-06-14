"""
Builds a professional PowerPoint slide deck from YouTube AI insights.

Input:  .tmp/insights.json + .tmp/charts/*.png
Output: .tmp/ai_insights_report.pptx
"""

import json
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# Colors
C_BG = RGBColor(0x1A, 0x1A, 0x2E)
C_ACCENT1 = RGBColor(0xE9, 0x45, 0x60)
C_ACCENT2 = RGBColor(0x0F, 0x34, 0x60)
C_ACCENT3 = RGBColor(0x16, 0x21, 0x3E)
C_TEXT = RGBColor(0xEA, 0xEA, 0xEA)
C_MUTED = RGBColor(0xA0, 0xA0, 0xC0)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

CHART_DIR = Path(".tmp/charts")
OUT_PATH = Path(".tmp/ai_insights_report.pptx")
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def fmt_num(n) -> str:
    n = int(n or 0)
    if n >= 1_000_000:
        return f"{n/1_000_000:.1f}M"
    if n >= 1_000:
        return f"{n/1_000:.0f}K"
    return str(n)


def set_bg(slide, prs):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = C_BG


def add_textbox(slide, text, left, top, width, height,
                font_size=12, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or C_TEXT
    return txBox


def add_accent_bar(slide, top=Inches(0.55), height=Inches(0.06)):
    bar = slide.shapes.add_shape(1, 0, top, SLIDE_W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ACCENT1
    bar.line.fill.background()


def add_header(slide, title: str, subtitle: str = ""):
    add_accent_bar(slide)
    add_textbox(slide, title, Inches(0.4), Inches(0.08),
                Inches(12), Inches(0.5), font_size=22, bold=True, color=C_WHITE)
    if subtitle:
        add_textbox(slide, subtitle, Inches(0.4), Inches(0.7),
                    Inches(12), Inches(0.3), font_size=11, color=C_MUTED, italic=True)


def add_kpi_box(slide, left, top, width, height, label: str, value: str, value_color=None):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = C_ACCENT3
    box.line.color.rgb = C_ACCENT2

    tf = box.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = value
    r1.font.size = Pt(26)
    r1.font.bold = True
    r1.font.color.rgb = value_color or C_ACCENT1

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(10)
    r2.font.color.rgb = C_MUTED


def insert_chart(slide, filename: str, left, top, width, height):
    path = CHART_DIR / filename
    if path.exists():
        slide.shapes.add_picture(str(path), left, top, width, height)
    else:
        placeholder = slide.shapes.add_shape(1, left, top, width, height)
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = C_ACCENT3
        tf = placeholder.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"[{filename} not found]"
        r.font.color.rgb = C_MUTED
        r.font.size = Pt(10)


def add_bullet_list(slide, items: list[str], left, top, width, height, font_size=11):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = f"• {item}"
        r.font.size = Pt(font_size)
        r.font.color.rgb = C_TEXT


def slide_title(prs, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs)

    month_year = datetime.now().strftime("%B %Y")

    # Big gradient-ish title area
    hero = slide.shapes.add_shape(1, 0, Inches(1.8), SLIDE_W, Inches(3.2))
    hero.fill.solid()
    hero.fill.fore_color.rgb = C_ACCENT2
    hero.line.fill.background()

    add_textbox(slide, "📊 AI YouTube Insights Report",
                Inches(0.5), Inches(2.0), Inches(12), Inches(1.0),
                font_size=38, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, month_year,
                Inches(0.5), Inches(3.1), Inches(12), Inches(0.6),
                font_size=20, color=C_ACCENT1, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Trends · Engagement · Strategy · Opportunities",
                Inches(0.5), Inches(3.8), Inches(12), Inches(0.5),
                font_size=13, color=C_MUTED, align=PP_ALIGN.CENTER, italic=True)

    add_accent_bar(slide, top=Inches(1.75), height=Inches(0.08))
    add_accent_bar(slide, top=Inches(5.1), height=Inches(0.06))

    add_textbox(slide, "Powered by Apify · YouTube Data API · Claude AI",
                Inches(0.5), Inches(5.3), Inches(12), Inches(0.4),
                font_size=9, color=C_MUTED, align=PP_ALIGN.CENTER)


def slide_executive_summary(prs, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs)
    add_header(slide, "Executive Summary", "Key metrics from this report")

    s = data["summary"]
    kpis = [
        ("Videos Analyzed", fmt_num(s["total_videos_analyzed"])),
        ("Avg Engagement Rate", f"{s['avg_engagement_rate']:.2f}%"),
        ("Top Creator", s["top_creator"][:18] + "…" if len(s["top_creator"]) > 18 else s["top_creator"]),
        ("Total Views Tracked", fmt_num(s["total_views_across_all"])),
    ]

    kpi_w = Inches(2.9)
    kpi_h = Inches(1.8)
    kpi_top = Inches(1.5)
    for i, (label, value) in enumerate(kpis):
        left = Inches(0.35 + i * 3.1)
        add_kpi_box(slide, left, kpi_top, kpi_w, kpi_h, label, value)

    add_textbox(slide, "🏆 Top Trending Video",
                Inches(0.4), Inches(3.6), Inches(3), Inches(0.4),
                font_size=11, bold=True, color=C_ACCENT1)
    add_textbox(slide, s.get("top_video_title", "N/A"),
                Inches(0.4), Inches(4.0), Inches(12.5), Inches(0.5),
                font_size=14, bold=True, color=C_WHITE)
    add_textbox(slide, f"Views: {fmt_num(s.get('top_video_views', 0))}",
                Inches(0.4), Inches(4.55), Inches(4), Inches(0.35),
                font_size=11, color=C_MUTED)

    # Queries used
    queries = data.get("queries_used", [])
    if queries:
        add_textbox(slide, "Search queries: " + " · ".join(f'"{q}"' for q in queries[:5]),
                    Inches(0.4), Inches(5.5), Inches(12.5), Inches(0.5),
                    font_size=9, color=C_MUTED, italic=True)


def slide_top_views(prs, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs)
    add_header(slide, "Top Trending Videos", "Highest view counts in the last 30 days")
    insert_chart(slide, "top_videos_views.png", Inches(0.3), Inches(0.9), Inches(12.7), Inches(6.3))


def slide_top_engagement(prs, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs)
    add_header(slide, "Highest Engagement Videos", "Most reactions relative to view count (likes + comments / views)")
    insert_chart(slide, "top_videos_engagement.png", Inches(0.3), Inches(0.9), Inches(12.7), Inches(6.3))


def slide_channel_landscape(prs, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs)
    add_header(slide, "Channel Landscape", "Top AI channels ranked by engagement and average views")
    insert_chart(slide, "channel_comparison.png", Inches(0.3), Inches(0.9), Inches(12.7), Inches(6.3))


def slide_title_keywords(prs, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs)
    add_header(slide, "What's Working: Title Keywords", "Most frequent words in trending AI video titles")
    insert_chart(slide, "title_keywords.png", Inches(0.3), Inches(0.9), Inches(12.7), Inches(6.3))


def slide_upload_timing(prs, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs)
    add_header(slide, "Best Time to Post", "Upload activity patterns by day and hour (UTC)")
    insert_chart(slide, "upload_heatmap.png", Inches(0.3), Inches(0.9), Inches(12.7), Inches(6.3))


def slide_rising_stars(prs, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs)
    add_header(slide, "Rising Stars", "Videos gaining views fastest — high velocity = algorithm momentum")
    insert_chart(slide, "view_velocity.png", Inches(0.3), Inches(0.9), Inches(12.7), Inches(6.3))


def slide_recommendations(prs, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs)
    add_header(slide, "Strategic Recommendations", "Action items based on this week's data")

    s = data["summary"]
    top_keywords = [k["word"] for k in data.get("keywords", [])[:5]]
    top_channel = s.get("top_creator", "N/A")
    avg_eng = s.get("avg_engagement_rate", 0)

    # Find best upload day
    by_day = data.get("upload_by_day", {})
    best_day = max(by_day, key=lambda d: by_day[d]) if by_day else "Tuesday"

    recs = [
        f"Focus titles on high-frequency terms: {', '.join(top_keywords)} — these words dominate trending content right now.",
        f"Study {top_channel}'s content format and posting cadence — they lead in engagement this period.",
        f"Aim for >{avg_eng:.1f}% engagement rate; videos below this threshold are underperforming vs the niche average.",
        f"Post on {best_day}s — that's when AI content sees peak upload and viewership activity.",
        "Prioritize view velocity in first 48 hours: strong early engagement is the strongest signal to the YouTube algorithm.",
    ]

    for i, rec in enumerate(recs):
        top_pos = Inches(1.2 + i * 1.1)
        num_box = slide.shapes.add_shape(1, Inches(0.35), top_pos, Inches(0.45), Inches(0.45))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = C_ACCENT1
        num_box.line.fill.background()
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(i + 1)
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = C_WHITE

        add_textbox(slide, rec, Inches(0.9), top_pos, Inches(12.1), Inches(0.65),
                    font_size=11, color=C_TEXT)


def main():
    data = json.loads(Path(".tmp/insights.json").read_text())
    # Inject queries list if missing
    if "queries_used" not in data:
        data["queries_used"] = [
            "AI automation 2025", "AI tools", "artificial intelligence tutorial",
            "ChatGPT", "Claude AI", "LLM", "prompt engineering", "AI agents",
        ]

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building slides...")
    slide_title(prs, data)
    print("  ✓ Slide 1: Title")
    slide_executive_summary(prs, data)
    print("  ✓ Slide 2: Executive Summary")
    slide_top_views(prs, data)
    print("  ✓ Slide 3: Top Trending Videos")
    slide_top_engagement(prs, data)
    print("  ✓ Slide 4: Top Engagement")
    slide_channel_landscape(prs, data)
    print("  ✓ Slide 5: Channel Landscape")
    slide_title_keywords(prs, data)
    print("  ✓ Slide 6: Title Keywords")
    slide_upload_timing(prs, data)
    print("  ✓ Slide 7: Upload Timing")
    slide_rising_stars(prs, data)
    print("  ✓ Slide 8: Rising Stars")
    slide_recommendations(prs, data)
    print("  ✓ Slide 9: Recommendations")

    prs.save(str(OUT_PATH))
    print(f"\nPresentation saved → {OUT_PATH}")


if __name__ == "__main__":
    main()
